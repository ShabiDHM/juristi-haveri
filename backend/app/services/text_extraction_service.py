# FILE: backend/app/services/text_extraction_service.py
# PHOENIX PROTOCOL - TEXT EXTRACTION V8.5 (ACCOUNTING TERMINOLOGY ALIGNMENT)
# 1. REFACTOR: FOOTER_PATTERN updated to strip "Rasti:", "Biznesi:", and "Klienti:" footers.
# 2. FIX: Replaced hardcoded /tmp/ with tempfile logic for cross-platform stability.
# 3. STATUS: Optimized for Accounting (XLSX, CSV, PDF) ingestion. 100% Aligned.

import fitz  # PyMuPDF
import docx
from pptx import Presentation
import pandas as pd
import csv
from typing import Dict, Callable, Any, Optional, Union
import logging
import cv2
import numpy as np
import io
import concurrent.futures
import time
import os
import tempfile
import uuid
import re

# PHOENIX: Absolute-style import logic for consistency
try:
    from app.services.ocr_service import extract_text_from_image as advanced_image_ocr
except ImportError:
    advanced_image_ocr = None

logger = logging.getLogger(__name__)

MAX_WORKERS = 2
# Updated to catch all variations of client/business labels in the footer
FOOTER_PATTERN = re.compile(r'(Rasti|Biznesi|Klienti):\s*\S+\s*\|\s*Kontabilisti AI System')

def _sanitize_text(text: str) -> str:
    """Removes null bytes and cleans text."""
    if not text: return ""
    return text.replace("\x00", "")

def _strip_footer(text: str) -> str:
    """Remove the branding footer line(s) from extracted text."""
    lines = text.split('\n')
    filtered = [line for line in lines if not FOOTER_PATTERN.search(line)]
    return '\n'.join(filtered)

def _sort_blocks(blocks):
    """Sort text blocks by vertical then horizontal position."""
    return sorted(blocks, key=lambda b: (int(b[1] / 3), int(b[0])))

def _process_single_page_safe(doc_path: str, page_num: int) -> str:
    """
    Extracts text from a single PDF page with Layout Preservation.
    Specifically tuned for fiscal documents (Invoices/Pasqyra).
    """
    page_marker = f"\n--- [FAQJA {page_num + 1}] ---\n"
    try:
        with fitz.open(doc_path) as doc:
            page: Any = doc[page_num]
            
            # 1. Try Direct Text Extraction
            blocks = page.get_text("blocks")
            text = ""
            if blocks:
                sorted_blocks = _sort_blocks(blocks)
                text = "\n".join([b[4] for b in sorted_blocks])
                text = _sanitize_text(text)
            
            # Remove accounting branding footer
            text_without_footer = _strip_footer(text)
            
            # If meaningful text remains, return it.
            if text_without_footer and len(text_without_footer.strip()) > 50:
                return page_marker + text_without_footer
            
            # 2. Otherwise, trigger OCR (likely a scanned invoice)
            logger.info(f"Faqja {page_num} duket e skanuar. Duke nisur OCR...")
            
            if not advanced_image_ocr:
                return page_marker + "[DOKUMENT I SKANUAR - OCR JO AKTIV]"

            # Render page to high-res image for OCR
            zoom = 2.0 
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat) 
            img_data = pix.tobytes("png")
            
            temp_img_path = os.path.join(tempfile.gettempdir(), f"page_{page_num}_{uuid.uuid4()}.png")
            with open(temp_img_path, "wb") as f:
                f.write(img_data)
            
            ocr_text = ""
            try:
                ocr_text = advanced_image_ocr(temp_img_path)
            except Exception as ocr_err:
                logger.error(f"OCR Dështoi për faqen {page_num}: {ocr_err}")
            finally:
                if os.path.exists(temp_img_path):
                    os.remove(temp_img_path)

            return page_marker + _sanitize_text(ocr_text)

    except Exception as e:
        logger.error(f"Page {page_num} CRITICAL FAILURE: {e}")
        return "" 

def _process_single_page_wrapper(args) -> str:
    try: return _process_single_page_safe(*args)
    except Exception: return ""

def _extract_text_sequentially(file_path: str, total_pages: int) -> str:
    buffer = []
    for i in range(total_pages):
        buffer.append(_process_single_page_safe(file_path, i))
        time.sleep(0.05) 
    return "".join(buffer)

def _extract_text_from_pdf(file_path: str) -> str:
    logger.info(f"🚀 Procesimi i PDF: {file_path}")
    total_pages = 0
    try:
        with fitz.open(file_path) as doc:
            total_pages = len(doc)
    except Exception: return ""

    # Sequential processing for short docs, parallel for long ones
    if total_pages < 5:
         return _extract_text_sequentially(file_path, total_pages)

    try:
        page_args = [(file_path, i) for i in range(total_pages)]
        results = []
        with concurrent.futures.ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
            futures = {executor.submit(_process_single_page_wrapper, arg): arg[1] for arg in page_args}
            for future in concurrent.futures.as_completed(futures):
                results.append((futures[future], future.result()))
        
        results.sort(key=lambda x: x[0])
        return "".join([r[1] for r in results])
    except Exception as e:
        logger.error(f"Parallel extraction failed: {e}")
        return _extract_text_sequentially(file_path, total_pages)

def _extract_text_from_docx(file_path: str) -> str:
    try:
        doc = docx.Document(file_path)
        return _sanitize_text("\n".join(para.text for para in doc.paragraphs))
    except Exception as e:
        logger.error(f"DOCX Error: {e}")
        return ""

def _extract_text_from_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return _sanitize_text("\n".join(text_runs))
    except Exception as e:
        logger.error(f"PPTX Error: {e}")
        return ""

def _extract_text_from_image(file_path: str) -> str:
    if advanced_image_ocr:
        return _sanitize_text(advanced_image_ocr(file_path))
    return ""

def _extract_text_from_txt(file_path: str) -> str:
    try: 
        with open(file_path, 'r', encoding='utf-8') as f: 
            return _sanitize_text(f.read())
    except Exception: return ""

def _extract_text_from_csv(file_path: str) -> str:
    """Extracts accounting data from CSV format."""
    try:
        all_rows = []
        with open(file_path, mode='r', newline='', encoding='utf-8-sig') as csvfile:
            reader = csv.reader(csvfile)
            for row in reader: all_rows.append(", ".join(filter(None, row)))
        return _sanitize_text("\n".join(all_rows))
    except Exception: return ""

def _extract_text_from_excel(file_path: str) -> str:
    """Highly optimized for financial sheets and bank statements."""
    try:
        xls = pd.ExcelFile(file_path)
        full_text = []
        for sheet_name in xls.sheet_names:
            df = xls.parse(sheet_name)
            if df.empty: continue
            full_text.append(f"--- Faqja (Sheet): {sheet_name} ---")
            full_text.append(df.to_string(index=False, na_rep=""))
        return _sanitize_text("\n".join(full_text))
    except Exception: return ""

EXTRACTION_MAP: Dict[str, Callable[[str], str]] = {
    "application/pdf": _extract_text_from_pdf,
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": _extract_text_from_docx,
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": _extract_text_from_pptx,
    "image/png": _extract_text_from_image, 
    "image/jpeg": _extract_text_from_image, 
    "image/tiff": _extract_text_from_image,
    "image/jpg": _extract_text_from_image,
    "text/plain": _extract_text_from_txt, 
    "text/csv": _extract_text_from_csv,
    "application/vnd.ms-excel": _extract_text_from_excel,
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": _extract_text_from_excel
}

def extract_text(file_path: str, mime_type: str) -> str:
    normalized_mime_type = mime_type.split(';')[0].strip().lower()
    extractor = EXTRACTION_MAP.get(normalized_mime_type)
    
    if not extractor:
        if file_path.lower().endswith(('.png', '.jpg', '.jpeg', '.tiff')):
            return _extract_text_from_image(file_path)
        return _extract_text_from_txt(file_path)
        
    return extractor(file_path)

def extract_text_from_file(file_obj: io.BytesIO, file_type: str = "PDF") -> str:
    extension_map = {
        "PDF": ".pdf", "DOCX": ".docx", "PPTX": ".pptx",
        "PNG": ".png", "JPG": ".jpg", "JPEG": ".jpeg",
        "TXT": ".txt", "CSV": ".csv", "XLSX": ".xlsx"
    }
    ext = extension_map.get(file_type.upper(), ".bin")
    
    mime_map = {
        "PDF": "application/pdf", "DOCX": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "PPTX": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "XLSX": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "CSV": "text/csv", "TXT": "text/plain",
        "PNG": "image/png", "JPG": "image/jpeg"
    }
    mime_type = mime_map.get(file_type.upper(), "application/octet-stream")

    with tempfile.NamedTemporaryFile(suffix=ext, delete=False) as tmp:
        tmp.write(file_obj.getvalue())
        tmp_path = tmp.name

    try:
        return extract_text(tmp_path, mime_type)
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)